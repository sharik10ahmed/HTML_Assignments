from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 1. Initialize Presentation
prs = Presentation()

# Premium Color Palette
BG_COLOR = RGBColor(8, 8, 8)          # Zinc-950 (Dark Background)
GOLD_COLOR = RGBColor(201, 168, 76)   # Gold Accent
WHITE_COLOR = RGBColor(249, 246, 240) # Off-White Text
GRAY_COLOR = RGBColor(138, 134, 128)  # Zinc-400 Subtext

# 2. Presentation Data Extracted from HTML
slides_data = [
    {
        "tag": "Slide 01 · Core Engine",
        "title": "From Legacy Monolith to Async, Event-Driven Legal Infrastructure",
        "bullets": [
            "🚀 Async Engine: FastAPI powering non-blocking IO via routers and @app.websocket.",
            "🗄️ Data Layer: PostgreSQL driven by Async SQLAlchemy and Alembic.",
            "⚡ Event Layer: NATS Message Broker decoupled via Taskiq background workers.",
            "🎨 Render Layer: Jinja2 Server-side templating for client-side optimization."
        ]
    },
    {
        "tag": "Slide 02 · Authentication",
        "title": "One Identity Gateway. Four Roles. Zero Ambiguity.",
        "bullets": [
            "👤 Client: Case Creation, Document Upload, Razorpay Payment, Real-time Tracking.",
            "⚖️ Lawyer: Case Workspace Ops, Advance Pipeline Stages, Secure Client Chat.",
            "🛡️ Staff: Document Review Queue, Flag Malicious Reports, Warning Assertions.",
            "👑 Superadmin: Global Platform Override, Unrestricted Mutations, Audit Logs.",
            "Token Flow: Login → JWT Issued → LocalStorage → Role-Gated Routers."
        ]
    },
    {
        "tag": "Slide 03 · Data Models",
        "title": "Two Profiles. One Platform. A Legal Marketplace Architecture.",
        "bullets": [
            "🏛️ accounts_baseuser: Single identity record for every user (PK, email, password_hash, role).",
            "👤 clients_clientprofile: Extends baseuser (marital_status, gender, mobile_number, warnings_count).",
            "⚖️ lawyers_lawyerprofile: Extends baseuser (bar_reg_number, specialization, consultation_fee, rating)."
        ]
    },
    {
        "tag": "Slide 04 · Verification Gate",
        "title": "The Verification Wall: How Rogue Lawyers Never Reach Clients",
        "bullets": [
            "STEP 1 (Submitted): Lawyer completes onboarding. Status saved as PENDING.",
            "STEP 2 (Under Review): AdminPanelProfile staff accesses isolated queue (gated by is_staff==True).",
            "STEP 3 (Adjudicated): APPROVED → Profile goes live. REJECTED → Lawyer notified for resubmission."
        ]
    },
    {
        "tag": "Slide 05 · Case Flow Part I",
        "title": "The Case Lifecycle: Initiation to Assignment in Four Steps",
        "bullets": [
            "1. CASE_CREATED: Client selects verified lawyer. CaseRequest spawned with custom_id.",
            "2. DOCUMENT_VERIFICATION: Client uploads required legal documents to database arrays.",
            "3. DOCUMENTS_VERIFIED: Admin verifies inputs. Platform locks user modifications.",
            "4. LAWYER_ASSIGNED: Live case workspace activated. WebSocket notifications fired."
        ]
    },
    {
        "tag": "Slide 06 · Case Flow Part II",
        "title": "The Execution Engine: Seven Legal Milestones",
        "bullets": [
            "M1: Petition Drafted → M2: Petition Filed → M3: First Motion → M4: Second Motion",
            "M5: Decree Issued → M6: Completed → M7: Rated",
            "System Cascade: DB Status Mutation → Audit Trail Logged → WebSocket Broadcast → Taskiq Email."
        ]
    },
    {
        "tag": "Slide 07 · Storage Engine",
        "title": "The Document Engine: Secure Upload & Persistent Storage",
        "bullets": [
            "📤 UPLOAD LAYER: Client submits validation data (e.g., [Aadhaar Redacted] Card, Marriage Cert).",
            "☁️ STORAGE LAYER: upload_to_cloudinary() pipes binary stream server-side to CDN.",
            "📋 METADATA LAYER: Tracking records instantiated in lawyers_casedocument table."
        ]
    },
    {
        "tag": "Slide 08 · Aggregation Logic",
        "title": "The Verification Gate: The Last Line of Defence",
        "bullets": [
            "State Evaluation: Every admin approval re-queries all case documents.",
            "If Doc A (Verified) + Doc B (Pending) → Case BLOCKED.",
            "If All Documents == VERIFIED → Case AUTO-ADVANCES to LAWYER_ASSIGNED.",
            "Self-Executing: Backend query fires after every admin action without manual trigger."
        ]
    },
    {
        "tag": "Slide 09 · Real-Time Network",
        "title": "Always-On: The Dual-Channel Notification Infrastructure",
        "bullets": [
            "⚡ IN-SESSION: Real-Time WebSocket Layer (<50ms delivery).",
            "Maintains dynamic execution states. Triggers: Stage Changes, Chat Hooks, Payments.",
            "📬 OUT-OF-SESSION: Async Background Layer (NATS via Taskiq).",
            "Offloads processes to localized message queues. Triggers: Account Warnings, Receipt Dispatch."
        ]
    },
    {
        "tag": "Slide 10 · Chat Messaging",
        "title": "Case Chat: Persistent, Real-Time, Attachment-Capable",
        "bullets": [
            "Data Model: lawyers_casemessage (case_id, sender_type, text, attachment).",
            "Dispatch Sequence prioritizes persistence over real-time delivery.",
            "POST /messages → DB INSERT committed → manager.send_personal_message() → UI updates."
        ]
    },
    {
        "tag": "Slide 11 · Ledger & Payments",
        "title": "The Payments Engine: Razorpay-Verified & Audit-Complete",
        "bullets": [
            "STAGE 1 (Initiation): Payment DB row created with status PENDING.",
            "STAGE 2 (Capture): PCI-compliant Razorpay modal launches. Payment captured.",
            "STAGE 3 (Verification): HMAC-SHA256 signature validated against Razorpay secret.",
            "If valid → DB status SUCCEEDED → WebSocket Alert → NATS Receipt Email."
        ]
    },
    {
        "tag": "Slide 12 · Governance",
        "title": "Trust & Safety at Scale: Automated Enforcement",
        "bullets": [
            "Track A (Manual Report): Client flags Lawyer → Admin reviews → Triggers WARN.",
            "Track B (Automated Rules): Lawyer warnings ≥ 3 → Auto-Ban (is_active=False).",
            "Track B (Automated Rules): False reports ≥ 6 → Reporter Account Deactivated.",
            "14-Day Retention: Deletion initiated → 14-Day Cooling Window → Permanent DB Row Purge."
        ]
    }
]

# 3. Slide Generator Loop
for slide_data in slides_data:
    # Use blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    
    # Apply Premium Dark Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Add Tag/Eyebrow (Gold)
    txBox_tag = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    tf_tag = txBox_tag.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = slide_data["tag"].upper()
    p_tag.font.color.rgb = GOLD_COLOR
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True

    # Add Title (White)
    txBox_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.5))
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = slide_data["title"]
    p_title.font.color.rgb = WHITE_COLOR
    p_title.font.size = Pt(32)
    p_title.font.bold = True

    # Add Body Content (Gray/White)
    txBox_body = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4.5))
    tf_body = txBox_body.text_frame
    tf_body.word_wrap = True
    
    for bullet in slide_data["bullets"]:
        p_body = tf_body.add_paragraph()
        p_body.text = bullet
        p_body.font.color.rgb = GRAY_COLOR
        p_body.font.size = Pt(18)
        p_body.space_after = Pt(14)
        p_body.level = 0

# 4. Save the File
file_name = "DivorceConnect_Premium_Deck.pptx"
prs.save(file_name)
print(f"Success! {file_name} has been generated.")